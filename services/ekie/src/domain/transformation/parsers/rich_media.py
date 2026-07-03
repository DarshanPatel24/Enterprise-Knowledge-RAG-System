"""Generic rich-media parser for office docs, PDFs, and standalone images.

This parser provides broad format coverage now, while specialized parser-per-type
implementations can be introduced later. Extraction strategy is progressive:
1. unstructured (if installed)
2. format-specific fallback libraries (pypdf, python-docx, python-pptx)
3. OCR fallback for images (Pillow + pytesseract)
"""

from __future__ import annotations

from typing import ClassVar

from domain.transformation.parsers.base import (
    DocumentParser,
    ParsedContent,
    ParserContext,
    ParserError,
)

_IMAGE_EXTENSIONS = {
    "png",
    "jpg",
    "jpeg",
    "gif",
    "bmp",
    "tif",
    "tiff",
    "webp",
}


class RichMediaParser(DocumentParser):
    """Parser with broad coverage for enterprise binary document formats."""

    parser_type: ClassVar[str] = "rich_media"
    parser_version: ClassVar[str] = "1.0"
    supported_extensions: ClassVar[frozenset[str]] = frozenset(
        {
            "pdf",
            "doc",
            "docx",
            "ppt",
            "pptx",
            *_IMAGE_EXTENSIONS,
        }
    )
    supported_mime_types: ClassVar[frozenset[str]] = frozenset(
        {
            "application/pdf",
            "application/msword",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/vnd.ms-powerpoint",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "image/png",
            "image/jpeg",
            "image/gif",
            "image/bmp",
            "image/tiff",
            "image/webp",
        }
    )

    def parse(self, data: bytes, context: ParserContext) -> ParsedContent:
        """Extract markdown body text from supported rich-media formats."""
        text = self._extract_with_unstructured(data, context)
        if not text:
            text = self._extract_with_format_fallback(data, context)
        if not text.strip():
            raise ParserError(
                f"unable to extract textual content from {context.source_path!r}; "
                "install extraction dependencies or provide a specialized parser"
            )
        return ParsedContent(body=text)

    def _extract_with_unstructured(self, data: bytes, context: ParserContext) -> str:
        """Try extraction using unstructured's auto partitioner."""
        try:
            from unstructured.partition.auto import partition
        except Exception:
            return ""
        try:
            elements = partition(
                file=data,
                file_filename=context.source_path,
                content_type=context.mime_type,
            )
            chunks = [
                str(getattr(element, "text", "")).strip()
                for element in elements
                if str(getattr(element, "text", "")).strip()
            ]
            return "\n\n".join(chunks)
        except Exception:
            return ""

    def _extract_with_format_fallback(self, data: bytes, context: ParserContext) -> str:
        """Use library-specific extraction fallback by extension."""
        ext = context.extension.lower()
        if ext == "pdf":
            # Try text extraction first; if empty (scanned PDF) fall back to OCR.
            text = self._extract_pdf(data)
            if not text.strip():
                text = self._extract_pdf_ocr(data)
            return text
        if ext in {"doc", "docx"}:
            return self._extract_docx(data)
        if ext in {"ppt", "pptx"}:
            return self._extract_pptx(data)
        if ext in _IMAGE_EXTENSIONS:
            return self._extract_image_ocr(data)
        return ""

    def _extract_pdf(self, data: bytes) -> str:
        try:
            from io import BytesIO

            from pypdf import PdfReader
        except Exception:
            return ""
        try:
            reader = PdfReader(BytesIO(data))
            pages = [page.extract_text() or "" for page in reader.pages]
            return "\n\n".join(chunk.strip() for chunk in pages if chunk.strip())
        except Exception:
            return ""

    def _extract_pdf_ocr(self, data: bytes) -> str:
        """OCR fallback for scanned/image-only PDFs using pymupdf + pytesseract.

        pymupdf renders pages natively without poppler so it works on Windows
        without any additional system binaries beyond Tesseract itself.
        """
        try:
            import fitz  # pymupdf
            import pytesseract
            from PIL import Image
        except Exception:
            return ""
        try:
            doc = fitz.open(stream=data, filetype="pdf")
            page_texts: list[str] = []
            for page in doc:
                # Render at 2x scale for better OCR quality.
                matrix = fitz.Matrix(2.0, 2.0)
                pixmap = page.get_pixmap(matrix=matrix)
                image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
                text = pytesseract.image_to_string(image).strip()
                if text:
                    page_texts.append(text)
            doc.close()
            return "\n\n".join(page_texts)
        except Exception:
            return ""

    def _extract_docx(self, data: bytes) -> str:
        try:
            from io import BytesIO

            from docx import Document as DocxDocument
        except Exception:
            return ""
        try:
            document = DocxDocument(BytesIO(data))
            lines = [paragraph.text.strip() for paragraph in document.paragraphs]
            return "\n\n".join(line for line in lines if line)
        except Exception:
            return ""

    def _extract_pptx(self, data: bytes) -> str:
        try:
            from io import BytesIO

            from pptx import Presentation
        except Exception:
            return ""
        try:
            presentation = Presentation(BytesIO(data))
            slides_text: list[str] = []
            for slide in presentation.slides:
                fragments: list[str] = []
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and text.strip():
                        fragments.append(text.strip())
                if fragments:
                    slides_text.append("\n".join(fragments))
            return "\n\n".join(slides_text)
        except Exception:
            return ""

    def _extract_image_ocr(self, data: bytes) -> str:
        try:
            from io import BytesIO

            import pytesseract
            from PIL import Image
        except Exception:
            return ""
        try:
            image = Image.open(BytesIO(data))
            text = pytesseract.image_to_string(image)
            return text.strip()
        except Exception:
            return ""
